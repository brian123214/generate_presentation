import os
import PIL
from PIL import Image
from langchain.llms import OpenAI
from langchain.chains import LLMChain
from langchain.utilities import GoogleSearchAPIWrapper
from langchain.prompts import PromptTemplate
from pptx import Presentation
from slides import add_slide
from images import create_image
from stability_sdk import client


# API stuff
os.environ["STABILITY_KEY"] = "YOUROWN"
os.environ["OPENAI_API_KEY"] = "YOUROWN"
stability_api = client.StabilityInference(
    key=os.environ['STABILITY_KEY'], # API Key reference.
    verbose=True, # Print debug messages.
    engine="stable-diffusion-xl-beta-v2-2-2", # Set the engine to use for generation.
    # Available engines: stable-diffusion-v1 stable-diffusion-v1-5 stable-diffusion-512-v2-0 stable-diffusion-768-v2-0
    # stable-diffusion-512-v2-1 stable-diffusion-768-v2-1 stable-diffusion-xl-beta-v2-2-2 stable-inpainting-v1-0 stable-inpainting-512-v2-0
)

for key, value in os.environ.items():
    print(f'{key}: {value}')

# Change background and task to your own
background = "The Bagel"\
            "I stopped to pick up the bagel"\
            "rolling away in the wind,"\
            "annoyed with myself"\
            "for having dropped it"\
            "as if it were a portent."\
            "Faster and faster it rolled,"\
            "with me running after it"\
            "bent low, gritting my teeth,"\
            "and I found myself doubled over"\
            "and rolling down the street"\
            "head over heels, one complete somersault"\
            "after another like a bagel"\
            "and strangely happy with myself."\
            "—David Ignatow"
task = 'Please read the assigned poem and create a presentation for the class. The presentation should take about 10-15 minutes and include the following: You must begin the presentation reading the poem aloud. Your slides must include but in no specific order: 1. an explanation of the title 2. identify an image and explain its significance to the poem 3. identify and explain three lines from the poem--image, meaning, 4.  provide biographical information about the author--current position, 5. engage the class with a question that they must answer in the written form- the question may be posed  6. include any images, videos, or sound that you think will enhance our 7. answer the question: Why did Billy Collins choose to include this poem in the collection?"'


# Intitial task
llm = OpenAI(temperature=0.9)
assignment_prompt = PromptTemplate(
    input_variables=["task"],
    template="I have a presentation assignment to do the following surrounded by quotes: '{task}'. Slide by slide, what should the content and title of each slide be? When saying a slide name, preface it with the word 'SLIDE' and when stating each bullet point for the contents for each slide, preface it with the word 'CONTENT'",
)
assignment = LLMChain(llm=llm, prompt=assignment_prompt)
answer = assignment.run(task)
answer = answer.split('\n')


# Parse answer
slide_content = dict()
cur_slide = ""
for line in answer:
    if 'SLIDE' in line:
        cur_slide = line
    if 'CONTENT' in line:
        # Reason why we replace ':' is because answer by llm is usually like "CONTENT: actual content"
        line = line.replace('CONTENT:', '')
        slide_content[cur_slide] = line


# Generate content and image description
create_content_prompt = PromptTemplate(
    input_variables=["background", "content"],
    template="Here is the background information:{background} Answer the follwowing: {content}",
)
create_content = LLMChain(llm=llm, prompt=create_content_prompt)

create_image_description_prompt = PromptTemplate(
    input_variables=["content"],
    template="Describe the significant image in the following: {content}",
)
create_image_description = LLMChain(llm=llm, prompt=create_image_description_prompt)


# Creating powerpoint
prs = Presentation()
image_num = 0
for title, content in slide_content.items():
    # Generate content for that specific requested content
    content = create_content.run({"background":background, "content":content}).replace('\n', '')

    # Create an image description of the generated content to make 
    image_description = create_image_description.run(content).replace('\n', '')
    print(image_description)

    # Create image from generated image description
    image = create_image(image_description, stability_api)
    image_name = str(image_num) + '.png'
    image.save(image_name)

    # Add slide
    prs = add_slide(title, image_name, content, prs)
    image_num += 1
prs.save('hax.pptx')