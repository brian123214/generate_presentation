from pptx.util import Inches

def add_slide(title_text, image_name, content_text, prs):
    # # If you want this is latout for title slide
    # slide_layout = prs.slide_layouts[0]
    # add regular slide layout

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # title
    title = slide.shapes.title
    title.text = title_text

    # image
    left = Inches(0.75)
    top = Inches(1.25)
    slide.shapes.add_picture(image_name, left, top, height=Inches(4.5))
    
    # text
    left = Inches(0.75)
    top = Inches(1.5) + Inches(4)
    width = Inches(9)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = content_text
    return prs